import io
import os
import json
import csv
import docx
import tempfile
import openpyxl

from pathlib import Path
from typing import List, Optional
from fastapi import HTTPException
from bs4 import BeautifulSoup

from llama_index.core import Document
from docling.document_converter import DocumentConverter, PdfFormatOption
from docling.datamodel.pipeline_options import PdfPipelineOptions, AcceleratorOptions, AcceleratorDevice
from docling.datamodel.base_models import InputFormat

from app.core.enums.file import FileExtension, EncodingType, ParsingConstants
from app.core.enums.http import ErrorMessage, HttpStatus
from app.core.logger import get_logger

logger = get_logger("extractor")

# Global singleton for DocumentConverter to avoid re-initializing models
_docling_converter: Optional[DocumentConverter] = None


def get_docling_converter() -> DocumentConverter:
  global _docling_converter
  if _docling_converter is None:
    try:
      import torch
      device = AcceleratorDevice.CUDA if torch.cuda.is_available() else AcceleratorDevice.CPU
      logger.info(f"Initializing DocumentConverter with device: {device}...")
    except ImportError:
      logger.warning("Torch not found, defaulting to CPU for Docling.")
      device = AcceleratorDevice.CPU

    pipeline_options = PdfPipelineOptions()
    pipeline_options.accelerator_options = AcceleratorOptions(
        num_threads=4, device=device
    )
    _docling_converter = DocumentConverter(
        format_options={
            InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)
        }
    )
  return _docling_converter


# -------------------------------------------------------------------------
# MAIN EXTRACTOR
# -------------------------------------------------------------------------
def extract_documents(file_bytes: bytes, filename: str, reader_map: dict, arg_map: dict) -> List[Document]:
  """
  Tiered Extraction Strategy:
  1. DoclingReader (Default for PDF, DOCX, PPTX, XLSX, MD, HTML, Images)
  2. LlamaIndex Specific Readers (Fallback)
  3. Manual Extractor (Last Resort)
  """
  ext = os.path.splitext(filename.lower())[1]

  # Skip Docling for simple structured sets like CSV/JSON as standard readers are better/faster
  use_docling = ext in [
      FileExtension.PDF, FileExtension.DOCX, FileExtension.PPTX,
      FileExtension.XLSX, FileExtension.MD, FileExtension.HTML,
      FileExtension.JPG, FileExtension.JPEG, FileExtension.PNG,
      FileExtension.BMP, FileExtension.TIFF
  ]

  with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
    tmp.write(file_bytes)
    tmp_path = tmp.name

  try:
    if use_docling:
      try:
        logger.info(f"Attempting Docling extraction for {filename}")
        converter = get_docling_converter()
        # Docling supports file path directly
        conv_res = converter.convert(tmp_path)
        md_text = conv_res.document.export_to_markdown()

        if md_text:
          # Create a single document with the full text
          docs = [Document(text=md_text, metadata={
              "file_name": filename,
              "file_type": ext,
              "processing_method": "docling",
              "reader_type": "DocumentConverter(CUDA)"
          })]
          os.remove(tmp_path)
          return docs
      except Exception as e:
        logger.warning(f"Docling failed for {filename}, falling back: {e}")

    # Tier 2: LlamaIndex Default Readers
    if ext in reader_map:
      try:
        logger.info(f"Attempting LlamaIndex reader fallback for {filename}")
        reader = reader_map[ext]
        arg_name = arg_map.get(ext, "file")
        docs = reader.load_data(**{arg_name: Path(tmp_path)})

        if docs:
          for doc in docs:
            doc.metadata.update({
                "file_name": filename,
                "file_type": ext,
                "processing_method": "llama_index_reader",
                "reader_type": type(reader).__name__
            })
          os.remove(tmp_path)
          return docs
      except Exception as e:
        logger.warning(f"LlamaIndex reader failed for {filename}: {e}")

    # Tier 3: Manual Fallback
    logger.info(f"Fallback to manual extraction for {filename}")
    text = _extract_text_manual(file_bytes, filename)
    os.remove(tmp_path)
    return [Document(text=text, metadata={"file_name": filename, "processing_method": "manual"})]

  except Exception as e:
    if os.path.exists(tmp_path):
      os.remove(tmp_path)
    logger.error(f"Total extraction failure for {filename}: {e}")
    raise e


# -------------------------------------------------------------------------
# FALLBACK EXTRACTORS
# -------------------------------------------------------------------------
def _extract_text_manual(file_bytes: bytes, filename: str) -> str:
  """
  Dispatch manual extraction function based on file extension.
  """
  ext = os.path.splitext(filename.lower())[1]
  match ext:
    case FileExtension.TXT | FileExtension.MD:
      return _extract_text_from_txt(file_bytes)
    case FileExtension.DOCX:
      return _extract_text_from_docx(file_bytes)
    case FileExtension.CSV:
      return _extract_text_from_csv(file_bytes)
    case FileExtension.JSON:
      return _extract_text_from_json(file_bytes)
    case FileExtension.PPTX:
      return _extract_text_from_pptx(file_bytes)
    case FileExtension.PDF:
      return _extract_text_from_pdf(file_bytes)
    case FileExtension.XLSX:
      return _extract_text_from_xlsx(file_bytes)
    case FileExtension.HTML:
      return _extract_text_from_html(file_bytes)
    case _:
      raise HTTPException(
          status_code=HttpStatus.UNPROCESSABLE_ENTITY.value,
          detail=f"{ErrorMessage.UNSUPPORTED_FILE_TYPE}: {ext}"
      )


# -------------------------------------------------------------------------
# SPECIFIC MANUAL EXTRACTORS
# -------------------------------------------------------------------------
def _extract_text_from_txt(binary_txt: bytes) -> str:
  """
  Try decoding a text file using multiple encodings.
  """
  for enc in [
      EncodingType.UTF8,
      EncodingType.UTF8_SIG,
      EncodingType.LATIN1,
      EncodingType.CP1252
  ]:
    try:
      return binary_txt.decode(enc)
    except UnicodeDecodeError:
      continue
  raise HTTPException(
      status_code=HttpStatus.UNPROCESSABLE_ENTITY.value,
      detail=ErrorMessage.UNABLE_TO_DECODE_TEXT
  )


def _extract_text_from_docx(binary_docx: bytes) -> str:
  """
  Extract all text from a DOCX file including table content.
  """
  doc = docx.Document(io.BytesIO(binary_docx))
  text_parts = [p.text for p in doc.paragraphs if p.text.strip()]

  for table in doc.tables:
    for row in table.rows:
      row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
      if row_text:
        text_parts.append(ParsingConstants.COLUMN_SEPARATOR.join(row_text))

  return '\n'.join(text_parts)


def _extract_text_from_csv(binary_csv: bytes) -> str:
  """
  Extract structured text from a CSV file.
  """
  text_content = binary_csv.decode(EncodingType.UTF8_SIG)
  csv_reader = csv.reader(io.StringIO(text_content))
  text_parts = []
  headers = None

  for i, row in enumerate(csv_reader):
    if i == 0:
      headers = row
      text_parts.append(ParsingConstants.HEADERS_PREFIX +
                        ParsingConstants.COLUMN_SEPARATOR.join(headers))
    else:
      if headers and len(row) == len(headers):
        row_dict = dict(zip(headers, row))
        row_text = ParsingConstants.COLUMN_SEPARATOR.join([
            f"{k}{ParsingConstants.KEY_VALUE_SEPARATOR}{v}" for k, v in row_dict.items() if v.strip()
        ])
      else:
        row_text = ParsingConstants.COLUMN_SEPARATOR.join(
            [cell for cell in row if cell.strip()])
      text_parts.append(row_text)

  return '\n'.join(text_parts)


def _extract_text_from_json(binary_json: bytes) -> str:
  """
  Extract structured text from a nested JSON file.
  """
  text_content = binary_json.decode(EncodingType.UTF8_SIG)
  data = json.loads(text_content)

  def extract_text(obj, prefix="") -> List[str]:
    parts = []
    if isinstance(obj, dict):
      for k, v in obj.items():
        current = f"{prefix}.{k}" if prefix else k
        parts.extend(extract_text(v, current))
    elif isinstance(obj, list):
      for i, item in enumerate(obj):
        current = f"{prefix}[{i}]" if prefix else f"item[{i}]"
        parts.extend(extract_text(item, current))
    else:
      parts.append(f"{prefix}: {str(obj)}")
    return parts

  return '\n'.join(extract_text(data))


def _extract_text_from_pptx(binary_pptx: bytes) -> str:
  """
  Extract text from a PPTX file including slides and notes.
  """
  from pptx import Presentation

  prs = Presentation(io.BytesIO(binary_pptx))
  text_parts = []

  for slide in prs.slides:
    # Extract text from shapes
    for shape in slide.shapes:
      if hasattr(shape, "text"):
        text_parts.append(shape.text)

    # Extract text from notes
    if slide.has_notes_slide:
      notes_slide = slide.notes_slide
      text_frame = notes_slide.notes_text_frame
      if text_frame:
        text_parts.append(text_frame.text)

  return '\n'.join(text_parts)


def _extract_text_from_pdf(binary_pdf: bytes) -> str:
  """
  Extract text from a PDF file using pymupdf4llm (Markdown).
  """
  import pymupdf4llm
  import tempfile

  with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
    tmp.write(binary_pdf)
    tmp_path = tmp.name

  try:
    return pymupdf4llm.to_markdown(tmp_path)
  except Exception as e:
    raise e
  finally:
    if os.path.exists(tmp_path):
      os.remove(tmp_path)


def _extract_text_from_xlsx(binary_xlsx: bytes) -> str:
  """
  Extract text from an XLSX file using openpyxl.
  """
  wb = openpyxl.load_workbook(io.BytesIO(binary_xlsx), data_only=True)
  text_parts = []
  for sheet in wb.worksheets:
    text_parts.append(f"Sheet: {sheet.title}")
    for row in sheet.iter_rows(values_only=True):
      row_text = [str(cell).strip() for cell in row if cell is not None]
      if row_text:
        text_parts.append(ParsingConstants.COLUMN_SEPARATOR.join(row_text))
  return '\n'.join(text_parts)


def _extract_text_from_html(binary_html: bytes) -> str:
  """
  Extract text from HTML using BeautifulSoup.
  """
  soup = BeautifulSoup(binary_html, "html.parser")
  # Remove script and style elements
  for script in soup(["script", "style"]):
    script.decompose()
  return soup.get_text(separator='\n', strip=True)
